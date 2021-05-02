from operator import itemgetter
import fitz
import json
import summary_gen as sg
from pptx.util import Pt
import os
from gtts import gTTS 
import re
from pptx import Presentation,util
from langdetect import detect 
import create_dict as cd

unacceptable = ['contents','content','table of content','references','table of contents','bibliography','acknowledgement','list of tables','list']
lang=""


class Node:
    def __init__(self, heading_tag, content, pg_no, parent, audio_link):
        self.heading = heading_tag
        self.content = content
        self.pg = pg_no
        self.parent = parent
        self.audio_link = audio_link
        self.children = []

def lang_detect(text):
  #print("text -- "+text)
  if text:
    return detect(text)
  else:
    return "en"

def fonts(doc, granularity=False):
    """Extracts fonts and their usage in PDF documents.
    :param doc: PDF document to iterate through
    :type doc: <class 'fitz.fitz.Document'>
    :param granularity: also use 'font', 'flags' and 'color' to discriminate text
    :type granularity: bool
    :rtype: [(font_size, count), (font_size, count}], dict
    :return: most used fonts sorted by count, font style information
    """
    styles = {}
    font_counts = {}

    for page in doc:
        blocks = page.getText("dict")["blocks"]
        for b in blocks:  # iterate through the text blocks
            if b['type'] == 0:  # block contains text
                for l in b["lines"]:  # iterate through the text lines
                    for s in l["spans"]:  # iterate through the text spans
                        s['size']=round(s['size'])
                        if granularity:
                            identifier = "{0}_{1}_{2}_{3}".format(s['size'], s['flags'], s['font'], s['color'])
                            styles[identifier] = {'size': s['size'], 'flags': s['flags'], 'font': s['font'],
                                                  'color': s['color']}
                        else:
                            identifier = "{0}".format(s['size'])
                            styles[identifier] = {'size': s['size'], 'font': s['font']}

                        font_counts[identifier] = font_counts.get(identifier, 0) + 1  # count the fonts usage

    font_counts = sorted(font_counts.items(), key=itemgetter(1), reverse=True)

    if len(font_counts) < 1:
        raise ValueError("Zero discriminating fonts found!")

    return font_counts, styles

def font_tags(font_counts, styles):
    """Returns dictionary with font sizes as keys and tags as value.
    :param font_counts: (font_size, count) for all fonts occuring in document
    :type font_counts: list
    :param styles: all styles found in the document
    :type styles: dict
    :rtype: dict
    :return: all element tags based on font-sizes
    """
    p_style = styles[font_counts[0][0]]  # get style for most used font by count (paragraph)
    p_size = p_style['size']  # get the paragraph's size

    # sorting the font sizes high to low, so that we can append the right integer to each tag
    font_sizes = []
    for (font_size, count) in font_counts:
        font_sizes.append(float(font_size))
    font_sizes.sort(reverse=True)

    main_heading = None

    # aggregating the tags for each font size
    idx = 0
    size_tag = {}
    freq_tag = {}
    for size in font_sizes:
        idx += 1
        if size == p_size:
            idx = 0
            size_tag[size] = '<p>'
        if size > p_size:
            size_tag[size] = '<h{0}>'.format(idx)
        elif size < p_size:
            size_tag[size] = '<s{0}>'.format(idx)
    
    font_counts = sorted(font_counts, key=lambda x:x[0], reverse=True)
    if font_counts[0][1]!=1:
        main_heading = size_tag[float(font_counts[0][0])]
    else:
        main_heading = size_tag[float(font_counts[1][0])]
    
    return size_tag, main_heading

def precedence_dict(size_tag):
  x=sorted(size_tag.items(), key=lambda x:x[0],reverse=True)
  d={}
  index_p=None
  for i in range(len(size_tag)):
    if x[i][1][1]=='p':
      index_p=i
      d[x[i][1]]=index_p
    elif x[i][1][1]=='s':
      d[x[i][1]]=index_p
    else:
      d[x[i][1]]=i
  d['<img>']=index_p
  return d

def get_images(pgno,pic_location,doc):
  list_images=[]
  for imgno,img in enumerate(doc.getPageImageList(pgno)):
      print("img found in ",pgno)
      xref = img[0]
      pix = fitz.Pixmap(doc, xref)
      if pix.n < 5:       # this is GRAY or RGB
          pix.writePNG(pic_location+"p%s-%s.png" % (pgno, imgno))
      else:               # CMYK: convert to RGB first
          pix1 = fitz.Pixmap(fitz.csRGB, pix)
          pix1.writePNG(pic_location+"p%s-%s.png" % (pgno, imgno))
          pix1 = None
      pix = None
      img_name= "p"+str(pgno)+"-"+str(imgno)+".png"
      path=os.path.abspath(pic_location+img_name)
      list_images.append("<img>"+path)
  return list_images

def check_figure(block_string):
  text=get_content(block_string)
  tag=get_tag(block_string)
  if re.match(r'^(f|F)ig*',text):
    return True
  return False

def get_string(text):
  x=str(text)
  text=x[2:-1]
  return text

def headers_para(doc, size_tag, pic_location):
    """Scrapes headers & paragraphs from PDF and return texts with element tags.
    :param doc: PDF document to iterate through
    :type doc: <class 'fitz.fitz.Document'>
    :param size_tag: textual element tags for each size
    :type size_tag: dict
    :rtype: list
    :return: texts with pre-prended element tags
    """
    header_para = []  # list with headers and paragraphs
    first = True  # boolean operator for first header
    previous_s = {}  # previous span
    block_string = ""
    flag=0
    global lang
    lang=""
    f=open("blocks.txt", 'w')
    for pgno,page in enumerate(doc):
        list_images=get_images(pgno,pic_location,doc)
        blocks = page.getText("dict")["blocks"]
        f.write(str(blocks).encode('ascii','ignore').decode('utf-8'))
        f.write("\n")
        #print("Blocks --------- ",blocks)
        for b in blocks:  # iterate through the text blocks
            if b['type'] == 0:  # this block contains text

                # REMEMBER: multiple fonts and sizes are possible IN one block

                #block_string = ""  # text found in block
                for l in b["lines"]:  # iterate through the text lines
                    for cnt,s in enumerate(l["spans"]):  # iterate through the text spans
                        #print(s)
                        if s['text'].strip().isdigit():
                          continue
                        s['size']=round(s['size'])
                        #if s['text']==" ":
                        #  block_string+=" || "
                        if len(s['text'])>=1 and s['text'].strip():  # removing whitespaces:
                            if first:
                              if lang_detect(s['text'])=="en":
                                lang="en"
                            s['text']=s['text'].encode('ascii','ignore').decode('utf-8') if lang=="en" else s['text']
                            if first:
                                previous_s = s
                                first = False
                                block_string = size_tag[s['size']] + s['text']
                            else:
                                #s['text']=s['text'].strip()
                                if s['size'] == previous_s['size']:
                                    #print("#####here - ",previous_s['text'],"|",s['text'])
                                    #if block_string: #and all((c == "|") for c in block_string):
                                    # block_string only contains pipes
                                    #    block_string = size_tag[s['size']] + s['text']
                                    #if block_string == "":
                                    #    # new block has started, so append size tag
                                    #    block_string = size_tag[s['size']] + s['text']
                                    #else:  # in the same block, so concatenate strings
                                    if list_images and re.match(r'^(f|F)ig*',s['text'].strip()):
                                      img=list_images.pop(0)
                                      header_para.append(block_string)
                                      header_para.append(img)
                                      block_string = size_tag[s['size']] + s['text']
                                    else:
                                      block_string += " " + s['text']
                                else:
                                    #print(previous_s['text'],"|",previous_s['size'],"|",s['size'],"|",s['text'])
                                    if previous_s['text'] and previous_s['text'][-1]==" " or (s['text'] and s['text'][0]==" "):
                                      #print("flag")
                                      if header_para and not check_figure(header_para[-1]) and get_tag(block_string)==get_tag(header_para[-1]):
                                        header_para[-1] += " " + get_content(block_string)
                                      else:
                                        if list_images and check_figure(block_string):
                                            img=list_images.pop(0)
                                            header_para.append(img)
                                        header_para.append(block_string)
                                      block_string = size_tag[s['size']] + s['text']
                                    else:
                                      if cnt!=0:
                                        block_string += s['text']
                                      else:
                                        if header_para and not check_figure(header_para[-1]) and get_tag(block_string)==get_tag(header_para[-1]):
                                          header_para[-1] += " " + get_content(block_string)
                                        else:
                                          if list_images and check_figure(block_string):
                                            img=list_images.pop(0)
                                            header_para.append(img)
                                          header_para.append(block_string)
                                        block_string = size_tag[s['size']] + s['text']
                                if check_figure(block_string):
                                  print("**********",block_string)
                                  header_para.append(block_string)
                                  block_string=size_tag[previous_s['size']]
                                else:
                                  previous_s = s
                                #print(header_para)
                    #print("-------------")

                    # new block started, indicating with a pipe
                    # block_string += "|"
                    
                #header_para.append(block_string)
                #print("**************")
                  
        if list_images:
          header_para.extend(list_images)        
        #print("-------------------------------------")
    header_para.append(block_string)
    return header_para

def get_tag(text):
  span=regex_check(text)
  return text[span[0]:span[1]]

def get_content(text):
  span=regex_check(text)
  return text[span[1]:]

def regex_check(text):
  match=re.match(r'^\<(.*?)\>',text)
  return match.span()

def form_node(element,span,index):
  heading_tag=element[span[0]:span[1]]
  content=element[span[1]:]
  pg_no=index
  parent=''
  audio_link=''
  return Node(heading_tag,content,pg_no,parent,audio_link)

def filtering(elements,doc):
  freq_dict={}
  for cnt,i in enumerate(elements):
    text=get_content(i)
    tag=get_tag(i)
    text=text.strip()
    preprocessed_text=re.sub('  ','',re.sub('\|','',re.sub('\d', '', text)))
    final = tag+preprocessed_text
    if final in freq_dict:
      freq_dict[final][0]+=1
      freq_dict[final][1].append(cnt)
    else:
      freq_dict[final]=[1,[cnt],tag]
  x=sorted(freq_dict.items(),key = lambda x:x[1][0], reverse=True)
  #print("%%%%%%%%%%%%%%%%%%%%%%%%%%%%%%",x)
  upper_limit=len(doc)
  lower_limit=upper_limit*0.25
  indices_to_be_removed=[]
  for i in x[:10]:
    #print("**********************",i[1][-1],i[1][0],i[0],lower_limit)
    if (i[1][-1][1]=='s' and i[1][0]>=lower_limit):
      indices_to_be_removed.extend(i[1][1])
    elif(i[1][-1]=='<img>'):
      continue
    else:
      break
  indices_to_be_removed.sort(reverse=True)
  #print("^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^", indices_to_be_removed)
  for index in indices_to_be_removed:
    del elements[index]
  return elements

def traverse(node_obj,v,main,index):
  span=regex_check(v)
  tag=v[span[0]:span[1]]
  content=get_content(v)
  if ( not node_obj.children):
    if(tag[1]=="h" and content.lower().strip() not in unacceptable):
      index[0]+=1
    node_obj.children.insert(0,form_node(v,regex_check(v),index[0]))
    #print("tfirst",node_obj.heading)
    return
  else:
    current = node_obj.children
    for child in current:
      if (main[child.heading] >= main[tag]):
        if(tag[1]=="h" and content.lower().strip() not in unacceptable):
          index[0]+=1
        node_obj.children.insert(0,form_node(v,regex_check(v),index[0]))
        #print("tsecond",node_obj.heading)
        return
      else:
        if(child.content.lower().strip() not in unacceptable):
          traverse(child,v,main,index)
        return

def create_hierarchy(elements, size_tag):
    list_dll = []
    main = precedence_dict(size_tag)
    user_list = elements
    index=[-1]
    #print(main)
    current_node = None
    for v in user_list:
        span=regex_check(v)
        tag=v[span[0]:span[1]]
        heading_content = get_content(v)
        if(list_dll):
            if(main[list_dll[-1].heading] >= main[tag]):#h2#new heading
                if(tag[1]=='h' and heading_content.lower().strip() not in unacceptable):
                  index[0]+=1
                list_dll.append(form_node(v,regex_check(v),index[0]))
            else:
                if(list_dll[-1].content.lower().strip() not in unacceptable):
                  traverse(list_dll[-1],v,main,index)
        else:
            #print("third")
            if(tag[1]=='h' and heading_content.lower().strip() not in unacceptable):
              index[0]+=1
            list_dll.append(form_node(v,regex_check(v),index[0]))
        
    return list_dll

def preprocessing(filename, pic_location):

    document = filename
    doc = fitz.open(document)

    font_counts, styles = fonts(doc, granularity=False)
    #print(font_counts,styles)
    size_tag, main_heading = font_tags(font_counts, styles)
    #print(size_tag)
    elements = headers_para(doc, size_tag,pic_location)
    elements = filtering(elements,doc)
    #print(elements)
    with open("doc.json", 'w') as json_out:
        json.dump(elements, json_out)
    return elements, main_heading, size_tag

def remove_string_special_characters(s): 
        
        # removes special characters with ' ' 
        stripped = re.sub('[^a-zA-z\s]', '', s) 
        stripped = re.sub('_', '', stripped) 
        
        # Change any white space to one space 
        stripped = re.sub('\s+', ' ', stripped) 
        
        # Remove start and end white spaces 
        stripped = stripped.strip() 
        if stripped != '': 
                return stripped.lower()

def topic_gen(text):
    import nltk 
    import re 
    from sklearn.feature_extraction.text import CountVectorizer, TfidfVectorizer 
    from nltk.corpus import stopwords 
    from nltk.tokenize import word_tokenize 
    import pandas as pd 

    
    txt1=text.split(".")

    # Preprocessing 
            
    # Stopword removal 
    stop_words = set(stopwords.words('english')) 
    your_list = ['skills', 'ability', 'job', 'description'] 
    for i, line in enumerate(txt1): 
        txt1[i] = ' '.join([x for
            x in nltk.word_tokenize(line) if
            ( x not in stop_words ) and ( x not in your_list )])
        
    # Getting trigrams 
    vectorizer = CountVectorizer(ngram_range = (3,3)) 
    X1 = vectorizer.fit_transform(txt1) 
    features = (vectorizer.get_feature_names()) 
    # print("\n\nFeatures : \n", features) 
    # print("\n\nX1 : \n", X1.toarray()) 

    # Applying TFIDF 
    vectorizer = TfidfVectorizer(ngram_range = (3,3)) 
    X2 = vectorizer.fit_transform(txt1) 
    scores = (X2.toarray()) 
    # print("\n\nScores : \n", scores) 

    # Getting top ranking features 
    sums = X2.sum(axis = 0) 
    data1 = []
    features=list(features)
    for col, term in enumerate(features): 
        data1.append( (term, sums[0,col] )) 
    ranking = pd.DataFrame(data1, columns = ['term','rank']) 
    words3 = (ranking.sort_values('rank', ascending = False)).reset_index()
    print ("\n\nWords head : \n", words3.head(7))

    vectorizer = CountVectorizer(ngram_range = (2,2)) 
    X1 = vectorizer.fit_transform(txt1) 
    features = (vectorizer.get_feature_names()) 
    # print("\n\nFeatures : \n", features) 
    # print("\n\nX1 : \n", X1.toarray()) 

    # Applying TFIDF 
    vectorizer = TfidfVectorizer(ngram_range = (2,2)) 
    X2 = vectorizer.fit_transform(txt1) 
    scores = (X2.toarray()) 
    # print("\n\nScores : \n", scores) 

    # Getting top ranking features 
    sums = X2.sum(axis = 0) 
    data1 = []
    features=list(features)
    for col, term in enumerate(features): 
        data1.append( (term, sums[0,col] )) 
    ranking = pd.DataFrame(data1, columns = ['term','rank']) 
    words2 = (ranking.sort_values('rank', ascending = False)).reset_index()
    #print ("\n\nWords head : \n", words2.head(7))

    vectorizer = CountVectorizer(ngram_range = (1,1)) 
    X1 = vectorizer.fit_transform(txt1) 
    features = (vectorizer.get_feature_names()) 
    # print("\n\nFeatures : \n", features) 
    # print("\n\nX1 : \n", X1.toarray()) 

    # Applying TFIDF 
    vectorizer = TfidfVectorizer(ngram_range = (1,1)) 
    X2 = vectorizer.fit_transform(txt1) 
    scores = (X2.toarray()) 
    # print("\n\nScores : \n", scores) 

    # Getting top ranking features 
    sums = X2.sum(axis = 0) 
    data1 = []
    features=list(features)
    for col, term in enumerate(features): 
        data1.append( (term, sums[0,col] )) 
    ranking = pd.DataFrame(data1, columns = ['term','rank']) 
    words1 = (ranking.sort_values('rank', ascending = False)).reset_index()
    #print ("\n\nWords head : \n", words1.head(7))

    threshold=(words1['rank'][0]+words2['rank'][0]+words3['rank'][0])/3
    #print(words1['term'].loc[0],threshold)
    if words2['rank'][0]>=threshold:
        return words2['term'][0]
    elif words3['rank'][0]>=threshold:
        return words3['term'][0]
    else:
        return words1['term'][0]

def get_para_summary(element):
    #summary=sg.processing(eval('b\''+element.content+'\'').decode('utf-8'))[0]
    global lang
    summary=sg.processing(element.content,lang)[0]
    summary = summary.split('.')
    summary.pop()
    for i in range(len(summary)):
        if summary[i][1]=="\n":
            summary[i]=summary[i][2:]+"."
        else:
            summary[i]=summary[i][1:]+"."
    return summary

def add_heading_slide(prs, title_layout, title_slide_layout, element, main_heading):
    slide = prs.slides.add_slide(title_layout)
    title = slide.shapes.title
    title.text = element.content
    #voice_text=title.text
    #myobj = gTTS(text=voice_text, lang='en', slow=False) 
    #myobj.save("voiceover.mp3")
    title.text_frame.paragraphs[0].font.size=Pt(60)
    for child in element.children[::-1]:
        if child.content.strip().lower() in unacceptable:
            continue
        if child.heading[1]=='h':
            add_heading_slide(prs,title_layout,title_slide_layout, child, main_heading)
        else:
            if child.heading!='<img>':
              add_paragraph_slide(prs,title_slide_layout,child, element)

def set_gen(hierarchy, question_content, index):
  temp=''
  i=0
  step=0
  while(i<len(hierarchy.keys())):
    # if hierarchy[i]["children"][0]["content"]=="Contents:":
    #   i+=1
    #   continue
    temp+=hierarchy[i]["content"]
    for j in hierarchy[i]["children"]:
      if(j["heading"]=="<img>"):
        continue
      temp+=j["content"]
    if i%2==0:
      question_content[index]=temp
      temp=''
      index+=1
    i+=1

def audio_gen(page_mapping, voiceover_location):
  lang=lang_detect(page_mapping[0]["content"])
  temp=''
  for i in page_mapping:
    temp=''
    temp+=page_mapping[i]["content"]
    for j in page_mapping[i]["children"]:
      if(j["heading"]=="<img>"):
        continue
      temp+=j["content"]
    myobj = gTTS(text=temp, lang=lang, slow=False)
    path=voiceover_location+str(i)+".mp3"
    myobj.save(path)
    page_mapping[i]["audio_link"]=path
    


def add_paragraph_slide(prs, layout, element, parent):
        summary=get_para_summary(element)
        for i in range(0,len(summary),3):
            x=min(i+3,len(summary))
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            if i==0:
                title.text=parent.content
            else:
                #pg_cnt+=1
                title.text=parent.content+" contd."
            #voice_text=title.text+' '.join(summary[i:x])
            #myobj = gTTS(text=voice_text, lang='en', slow=False) 
            #myobj.save("voiceover.mp3")
            subtitle = slide.placeholders[1]
            # left = top = width = height = util.Inches(1.0)
            shapes = slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            # movie = slide.shapes.add_movie("voiceover.mp3", 
            #            left , top , width , height, 
            #            poster_frame_image=None, 
            #            mime_type='video/unknown')
            # tf.fit_text()# = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            title_shape = shapes.title
            for j in range(i,x):
                p = tf.add_paragraph()
                p.text = summary[j]
                p.level=0
            try:
              tf.fit_text(max_size=25)
            except:
              pass
            #print("^^^^^^^^^^^^^^^^")
            #print(tf.paragraphs[0].font.size)
            #print("^^^^^^^^^^^^^^^^")

        
def pptgen(hierarchy, main_heading, tid):
    unacceptable = ['contents','content','table of content','references','table of contents','bibliography','acknowledgement']
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[1]
    title_layout = prs.slide_layouts[0]
    slide_to_voice={}
    pg_cnt = 0
    for element in hierarchy:
        print("element -- ",element)
        if element.content.strip().lower() in unacceptable:
            continue
        if element.heading[1]=='h':
            add_heading_slide(prs, title_layout, title_slide_layout, element, main_heading)
        else:
            add_paragraph_slide(prs, title_slide_layout, element)
    ppt_path='static/downloads/'+tid+".pptx"
    prs.save(ppt_path)
        
def create_dict(hierarchy, flag):
  unacceptable = ['contents','content','table of content','references','table of contents','bibliography','acknowledgement']
  delete_list=[]
  x=len(hierarchy)
  for i in range(len(hierarchy)):
    #print(hierarchy[i].__dict__)
    if hierarchy[i].content.lower().strip() not in unacceptable: 
      hierarchy[i]=hierarchy[i].__dict__
      if hierarchy[i]['children']:
        print(hierarchy[i]["content"],hierarchy[i]["children"])
        children=create_dict(hierarchy[i]['children'],flag)
        hierarchy[i]['children']=children
    else:
      delete_list.append(i)
  for j in delete_list[::-1]:
    del hierarchy[j]
  for i in range(len(hierarchy)):
    if flag==1:
      hierarchy[i]["children"]=hierarchy[i]["children"][::-1]
    else:
      hierarchy[i]["children"]=hierarchy[i]["children"]
  return hierarchy

def generate_metadata(filename,tid,pic_location,voiceover_location):
    elements,main_heading,size_tag=preprocessing(filename,pic_location)
    print("*********** ",main_heading)
    hierarchy=create_hierarchy(elements,size_tag)
    #print(hierarchy)
    result={}
    cd.modified_dict(hierarchy,result,[0])
    for i in result:
      result[i]=create_dict([result[i]],0)[0]
      print("key-",i,", value-", result[i])
    # print("***********************")
    # print(create_dict(hierarchy,1))
    pptgen(hierarchy, main_heading, tid)
    audio_gen(result,voiceover_location)
    question_content={}
    set_gen(result,question_content,0)
    print("#######################",question_content)
    ppt_path='static/downloads/'+tid+".pptx"
    name=filename.split(".")[0]
    pdf_name=name+".pdf"
    ppt_path=os.path.abspath(ppt_path)
    pdf_path=ppt_path[:-4]+"pdf"
    hierarchy=create_dict(hierarchy,1)
    return {'ppt_path':ppt_path,'pdf_path':pdf_path,'mapping':hierarchy,'question_content':question_content,"modified_dict":result}

#generate_metadata('Report.pdf','5','./lol', './lol')